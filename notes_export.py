#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
notes_export.py - PowerPoint ノート抽出ツール

PowerPoint (.pptx) ファイルから発表者ノートを抽出し、
CSV、Markdown、JSON 形式で出力するツールです。

Usage:
    # Markdown形式で出力（区切り線付き）
    python notes_export.py slides.pptx -o notes.md --format md --md-separator="---"
    
    # CSV形式（空行区切り付き）
    python notes_export.py slides.pptx -o notes.csv --format csv --csv-blank-row
    
    # JSON形式
    python notes_export.py slides.pptx -o notes.json --format json
    
    # 標準出力へ
    python notes_export.py slides.pptx -o - --format md
"""
import argparse
import csv
import json
import sys
from pathlib import Path
from typing import List, Dict, Optional, Union, Any, TextIO

try:
    from pptx import Presentation
    from pptx.slide import Slide
except ImportError as e:
    print("ERROR: python-pptx がインストールされていません。", file=sys.stderr)
    print("  uv add python-pptx または pip install python-pptx", file=sys.stderr)
    raise


def get_slide_title(slide: Slide) -> str:
    """
    スライドのタイトルを取得する
    
    Args:
        slide: PowerPointスライドオブジェクト
        
    Returns:
        str: スライドのタイトル（見つからない場合は空文字列）
    """
    try:
        if slide.shapes.title and slide.shapes.title.text:
            return slide.shapes.title.text.strip()
    except Exception:
        pass
    
    # フォールバック: 最初のテキストを含むシェイプから取得
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and shape.text and shape.text.strip():
                return shape.text.strip().splitlines()[0]
        except Exception:
            continue
    
    return ""


def unique_preserve_order(seq: List[str]) -> List[str]:
    """
    重複を除去しながら順序を保持する
    
    Args:
        seq: 文字列のリスト
        
    Returns:
        List[str]: 重複を除去したリスト
    """
    seen = set()
    result = []
    for item in seq:
        if item not in seen:
            result.append(item)
            seen.add(item)
    return result


def get_slide_notes(slide: Slide) -> str:
    """
    スライドから発表者ノートを抽出する
    
    Args:
        slide: PowerPointスライドオブジェクト
        
    Returns:
        str: 抽出されたノートのテキスト
    """
    notes_texts = []
    
    try:
        notes_slide = slide.notes_slide
    except Exception:
        notes_slide = None
    
    if notes_slide is None:
        return ""
    
    # 1) ノート専用プレースホルダーから取得
    try:
        if getattr(notes_slide, "notes_text_frame", None) is not None:
            text = notes_slide.notes_text_frame.text or ""
            if text.strip():
                notes_texts.append(text.strip())
    except Exception:
        pass
    
    # 2) ノートページの他のテキストシェイプから収集
    try:
        for shape in notes_slide.shapes:
            try:
                if hasattr(shape, "has_text_frame") and shape.has_text_frame and shape.text_frame:
                    text = shape.text_frame.text or ""
                elif hasattr(shape, "text"):
                    text = shape.text or ""
                else:
                    text = ""
                
                text = text.strip()
                if text:
                    notes_texts.append(text)
            except Exception:
                continue
    except Exception:
        pass
    
    # 重複除去しつつ順序を保持
    notes_texts = unique_preserve_order(notes_texts)
    
    # ダブル改行で結合（読みやすさのため）
    return "\n\n".join(notes_texts)


def normalize_newlines(text: str) -> str:
    """
    改行コードを正規化する
    
    Args:
        text: 正規化対象のテキスト
        
    Returns:
        str: 改行コードが統一されたテキスト
    """
    return (text or "").replace("\r\n", "\n").replace("\r", "\n")


def export_csv(
    rows: List[Dict[str, Any]], 
    out_path: Union[str, Path], 
    encoding: str = "utf-8", 
    blank_row: bool = False
) -> None:
    """
    CSV形式でエクスポートする
    
    Args:
        rows: エクスポートするデータ
        out_path: 出力先パス（"-"の場合は標準出力）
        encoding: 文字エンコーディング
        blank_row: スライド間に空行を挿入するかどうか
    """
    if out_path == '-':
        file_handle: TextIO = sys.stdout
        should_close = False
    else:
        file_handle = Path(out_path).open("w", newline="", encoding=encoding)
        should_close = True
    
    try:
        writer = csv.writer(file_handle)
        writer.writerow(["slide_number", "title", "notes"])
        
        for i, row in enumerate(rows):
            writer.writerow([row["slide_number"], row["title"], row["notes"]])
            if blank_row and i < len(rows) - 1:
                writer.writerow([])  # 区切り用の空行
    finally:
        if should_close:
            file_handle.close()


def export_markdown(
    rows: List[Dict[str, Any]], 
    out_path: Union[str, Path], 
    encoding: str = "utf-8", 
    separator: Optional[str] = None, 
    pagebreak: bool = False
) -> None:
    """
    Markdown形式でエクスポートする
    
    Args:
        rows: エクスポートするデータ
        out_path: 出力先パス（"-"の場合は標準出力）
        encoding: 文字エンコーディング
        separator: スライド間の区切り文字（例: "---"）
        pagebreak: 印刷用改ページを挿入するかどうか
    """
    lines = ["# Speaker Notes Export\n"]
    
    for idx, row in enumerate(rows):
        title = row["title"] or "(無題)"
        lines.append(f"## Slide {row['slide_number']}: {title}")
        
        notes = (row["notes"] or "").strip()
        if notes:
            lines.append("")
            lines.append(notes)
            lines.append("")
        else:
            lines.append("\n_(ノートなし)_\n")
        
        # スライド間の区切り（最後のスライド以外）
        if idx < len(rows) - 1:
            if separator:
                lines.append(separator)
                lines.append("")  # スペーシング
            if pagebreak:
                # HTMLからPDF変換時の改ページ
                lines.append('<div style="page-break-after: always;"></div>')
                lines.append("")
    
    content = "\n".join(lines)
    
    if out_path == '-':
        sys.stdout.write(content)
    else:
        Path(out_path).write_text(content, encoding=encoding)


def export_json(
    rows: List[Dict[str, Any]], 
    out_path: Union[str, Path], 
    encoding: str = "utf-8"
) -> None:
    """
    JSON形式でエクスポートする
    
    Args:
        rows: エクスポートするデータ
        out_path: 出力先パス（"-"の場合は標準出力）
        encoding: 文字エンコーディング
    """
    content = json.dumps(rows, ensure_ascii=False, indent=2)
    
    if out_path == '-':
        sys.stdout.write(content)
    else:
        Path(out_path).write_text(content, encoding=encoding)


def parse_arguments() -> argparse.Namespace:
    """
    コマンドライン引数をパースする
    
    Returns:
        argparse.Namespace: パース結果
    """
    parser = argparse.ArgumentParser(
        description="PowerPoint (.pptx) から発表者ノートを抽出",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__
    )
    
    parser.add_argument(
        "pptx_path", 
        type=Path, 
        help="入力する .pptx ファイルのパス"
    )
    parser.add_argument(
        "-o", "--output", 
        type=Path, 
        help="出力ファイルのパス（省略時は入力ファイルと同じディレクトリ）"
    )
    parser.add_argument(
        "--format", 
        choices=["csv", "md", "json"], 
        default="csv", 
        help="出力形式（デフォルト: csv）"
    )
    parser.add_argument(
        "--encoding", 
        default="utf-8", 
        help="出力ファイルのエンコーディング（デフォルト: utf-8）"
    )
    
    # セパレータオプション
    parser.add_argument(
        "--csv-blank-row", 
        action="store_true", 
        help="CSV形式でスライド間に空行を挿入"
    )
    parser.add_argument(
        "--md-separator", 
        type=str, 
        default=None, 
        help="Markdown形式でスライド間に挿入する区切り文字（例: '---'）"
    )
    parser.add_argument(
        "--md-pagebreak", 
        action="store_true", 
        help="Markdown形式で印刷用改ページを挿入"
    )
    
    return parser.parse_args()


def main() -> None:
    """メイン処理"""
    args = parse_arguments()
    
    # 入力ファイルの存在確認
    if not args.pptx_path.exists():
        print(f"ERROR: 入力ファイルが見つかりません: {args.pptx_path}", file=sys.stderr)
        sys.exit(1)
    
    # プレゼンテーションを読み込み
    try:
        presentation = Presentation(str(args.pptx_path))
    except Exception as e:
        print(f"ERROR: PowerPointファイルの読み込みに失敗: {e}", file=sys.stderr)
        sys.exit(1)
    
    # 各スライドからデータを抽出
    rows: List[Dict[str, Any]] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        title = normalize_newlines(get_slide_title(slide))
        notes = normalize_newlines(get_slide_notes(slide))
        
        rows.append({
            "slide_number": idx,
            "title": title,
            "notes": notes,
        })
    
    # 出力パスの決定
    if args.output is not None:
        out_path = str(args.output)
    else:
        # デフォルトの出力パス
        stem = args.pptx_path.stem
        suffix_map = {"csv": ".csv", "md": ".md", "json": ".json"}
        suffix = suffix_map[args.format]
        out_path = str(args.pptx_path.with_name(f"{stem}_notes{suffix}"))
    
    # 形式に応じてエクスポート
    try:
        if args.format == "csv":
            export_csv(rows, out_path, args.encoding, blank_row=args.csv_blank_row)
        elif args.format == "md":
            export_markdown(
                rows, out_path, args.encoding, 
                separator=args.md_separator, 
                pagebreak=args.md_pagebreak
            )
        else:  # json
            export_json(rows, out_path, args.encoding)
        
        # 完了メッセージ（標準出力以外の場合）
        if out_path != '-':
            print(f"完了: {len(rows)} スライドのノートを抽出しました -> {out_path}")
    
    except Exception as e:
        print(f"ERROR: エクスポート中にエラーが発生: {e}", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()